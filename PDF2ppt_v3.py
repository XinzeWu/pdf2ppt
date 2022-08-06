#!/usr/bin/env python
# coding: utf-8

# # 每页2个的情况，可以制作为PPT

# In[1]:


# 加载库
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from pptx import Presentation
import cv2
import matplotlib.pyplot as plt
from tqdm import tqdm
import fitz
import re


# In[ ]:


def pdf2ppt(pdf_path,ppt_path,mode=1,slide = '16:9'):
    
    # 实例化 ppt 文档对象
    ppt = Presentation()
    # ＃ 设置幻灯片尺寸，16:9
    slide = slide.split(':')
    w = int(slide[0])
    h = int(slide[1])
    ppt.slide_width = Inches(w)
    ppt.slide_height = Inches(h)
    
    pdf_root = pdf_path
    pdf_path = os.listdir(pdf_root)
    for p in tqdm(pdf_path):
        path = f'{pdf_root}/{p}'
        try:
            pdf = fitz.Document(path)
        except:
            print(f'{p} is not a pdf')
        for i,pg in enumerate(range(0, len(pdf))):
            page = pdf[pg]  # 获得每一页的对象
            trans = fitz.Matrix(3.0, 3.0)#.preRotate(0)
            pm = page.get_pixmap(matrix=trans, alpha=False)  # 获得每一页的流对象
            pm.save('1.png')
            
            
            
            if mode == 1:
                              
                layout=ppt.slide_layouts[6] #空白布局
                slide=ppt.slides.add_slide(layout)
                #定义图片位置
                left=Inches(0)
                top=Inches(0)
                #定义图片大小
                width=Inches(16)
                height=Inches(9)
                img_path='1.png'
                pic=slide.shapes.add_picture(img_path,left,top,width,height)
            
            if mode == 2:
                
                img = cv2.imread('1.png')
                img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
                
                img1 = img[310:int(len(img)*0.44),180:-180]
                cv2.imwrite('1.png',img1[:,:,::-1],[int(cv2.IMWRITE_PNG_COMPRESSION), 0])
                layout=ppt.slide_layouts[6] #空白布局
                slide=ppt.slides.add_slide(layout)
                #定义图片位置
                left=Inches(0)
                top=Inches(0)
                #定义图片大小
                width=Inches(16)
                height=Inches(9)
                img_path='1.png'
                pic=slide.shapes.add_picture(img_path,left,top,width,height)
                
                img2 = img[int(len(img)*0.56):-310,180:-180]
                cv2.imwrite('1.png',img2[:,:,::-1],[int(cv2.IMWRITE_PNG_COMPRESSION), 0])
                layout=ppt.slide_layouts[6] #空白布局
                slide=ppt.slides.add_slide(layout)
                #定义图片位置
                left=Inches(0)
                top=Inches(0)
                #定义图片大小
                width=Inches(w)
                height=Inches(h)
                img_path='1.png'
                pic=slide.shapes.add_picture(img_path,left,top,width,height)


        ppt.save(f'ppt/{p[:-4]}.pptx')

    
    


# In[ ]:





# In[ ]:




