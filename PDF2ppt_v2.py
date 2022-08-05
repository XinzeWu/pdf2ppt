#!/usr/bin/env python
# coding: utf-8

# # 每页2个的情况，可以制作为PPT

# In[1]:


# 加载库
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from pptx import Presentation
import cv2
import matplotlib.pyplot as plt
from tqdm import tqdm
import fitz
import re


# In[21]:


# 实例化 ppt 文档对象
ppt = Presentation()
# ＃ 设置幻灯片尺寸，16:9
ppt.slide_width = Inches(16)
ppt.slide_height = Inches(9)

pdf_path = os.listdir('pdf')
for p in tqdm(pdf_path):
    path = f'pdf/{p}'
    pdf = fitz.Document(path)
    for i,pg in enumerate(range(0, len(pdf))):
        page = pdf[pg]  # 获得每一页的对象
        trans = fitz.Matrix(3.0, 3.0)#.preRotate(0)
        pm = page.get_pixmap(matrix=trans, alpha=False)  # 获得每一页的流对象
        pm.save('1.png')
        #读取图片，写入图片
        img = cv2.imread('1.png')
        img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
        img1 = img[310:int(len(img)*0.44),180:-180]
        cv2.imwrite('1.png',img1[:,:,::-1],[int(cv2.IMWRITE_PNG_COMPRESSION), 0])

        layout=ppt.slide_layouts[6] #空白布局
        slide=ppt.slides.add_slide(layout)
        #定义图片位置
        left=Inches(0)
        top=Inches(0)
        #定义图片大小
        width=Inches(16)
        height=Inches(9)
        img_path='1.png'
        pic=slide.shapes.add_picture(img_path,left,top,width,height)

        img2 = img[int(len(img)*0.56):-310,180:-180]
        cv2.imwrite('1.png',img2[:,:,::-1],[int(cv2.IMWRITE_PNG_COMPRESSION), 0])
        layout=ppt.slide_layouts[6] #空白布局
        slide=ppt.slides.add_slide(layout)
        #定义图片位置
        left=Inches(0)
        top=Inches(0)
        #定义图片大小
        width=Inches(16)
        height=Inches(9)
        img_path='1.png'
        pic=slide.shapes.add_picture(img_path,left,top,width,height)

    ppt.save(f'ppt/{p}.pptx')
    
    


# In[ ]:




